#!/usr/bin/env -S uv run --with python-pptx --script
"""
Deriva plantilla-irem.pptx del master institucional.

    ./derivar-plantilla-pptx.py ppt_resultados_IREM_2025_master_logos.pptx plantilla-irem.pptx

Esto NO se corre para hacer una presentacion: la plantilla ya viene en el
plugin. Se corre cuando la institucion publica un master nuevo, para volver a
derivarla. Conserva el slide master, los once layouts, el tema, los medios y
las cuatro variantes de Montserrat embebidas; quita las laminas de contenido y
los metadatos que arrastran sus titulos.
"""
import shutil
import sys
import zipfile
from pathlib import Path
from pptx import Presentation

APP_LIMPIO = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" \
xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">\
<Application>Microsoft Office PowerPoint</Application><PresentationFormat>Widescreen</PresentationFormat>\
<Slides>0</Slides><Notes>0</Notes><HiddenSlides>0</HiddenSlides><ScaleCrop>false</ScaleCrop>\
<Company>Iniciativa Regional para la Eliminacion de la Malaria</Company><LinksUpToDate>false</LinksUpToDate>\
<SharedDoc>false</SharedDoc><HyperlinksChanged>false</HyperlinksChanged>\
</Properties>"""


def main():
    if len(sys.argv) != 3:
        sys.exit(__doc__)
    origen, destino = Path(sys.argv[1]), Path(sys.argv[2])

    prs = Presentation(origen)
    lst = prs.slides._sldIdLst
    n = len(lst)
    # Al soltar la relacion, la lamina queda inalcanzable en el grafo del
    # paquete y python-pptx no la escribe. Con ella se van sus notas y los
    # medios que solo ella usaba, sin tener que tocar [Content_Types].xml.
    for sldId in list(lst):
        prs.part.drop_rel(sldId.rId)
        lst.remove(sldId)
    prs.core_properties.title = "Plantilla IREM / BID"
    prs.core_properties.author = "Iniciativa Regional para la Eliminacion de la Malaria"
    prs.core_properties.comments = ""
    tmp = destino.with_suffix(".tmp.pptx")
    prs.save(tmp)

    # docProps/app.xml sobrevive intacto a lo anterior y sigue listando los
    # titulos de las laminas que acabamos de quitar. Se reescribe entero.
    with zipfile.ZipFile(tmp) as z, zipfile.ZipFile(
        destino, "w", zipfile.ZIP_DEFLATED
    ) as salida:
        for item in z.infolist():
            datos = APP_LIMPIO.encode() if item.filename == "docProps/app.xml" else z.read(item.filename)
            salida.writestr(item, datos)
    tmp.unlink()

    v = Presentation(destino)
    print(f"quitadas {n} laminas; quedan {len(v.slides._sldIdLst)}")
    print(f"layouts ({len(v.slide_layouts)}): {', '.join(l.name for l in v.slide_layouts)}")
    with zipfile.ZipFile(destino) as z:
        fuentes = [i for i in z.namelist() if i.startswith("ppt/fonts/")]
        assert "docProps/app.xml" in z.namelist()
        app = z.read("docProps/app.xml").decode()
    assert "Slide Titles" not in app, "app.xml todavia arrastra titulos"
    print(f"fuentes embebidas: {len(fuentes)}   tamano: {destino.stat().st_size // 1024} KB")


if __name__ == "__main__":
    main()
