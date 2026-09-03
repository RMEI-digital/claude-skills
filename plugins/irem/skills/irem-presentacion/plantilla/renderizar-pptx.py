#!/usr/bin/env -S uv run --with python-pptx --with pillow --script
"""
Convierte una presentación IREM en PowerPoint editable.

    ./renderizar-pptx.py <archivo>.qmd     ->  <archivo>.pptx

El .qmd es el mismo del camino en PDF: una sola fuente, dos salidas. Lo que
cambia es el motor. Aquí no hay LaTeX: se abre `plantilla-irem.pptx`, que es el
master institucional sin sus láminas de contenido, y se van agregando láminas
sobre sus layouts. De ahí salen gratis el fondo con su degradado y su franja,
los logotipos del pie, la portada completa con la fotografía y la cintilla de
donantes, y la lámina de cierre.

Tres cosas que hay que saber antes de tocar este archivo:

1. LOS LAYOUTS MIENTEN. El master arrastra restos del export de Google Slides:
   el título de lámina viene declarado en 40 pt y en #0070C0, y el fontScheme
   del tema dice Arial. Nada de eso es la identidad. Por eso cada corrida fija
   a mano tamaño, color y tipografía en TODOS los renglones que escribe. Si
   alguna vez ves un título celeste o un texto en Arial, es que se escapó uno.

2. LA ESCALA NO ES LA DEL MASTER, ES LA DEL PDF. El master pide 14-16 pt de
   cuerpo sobre este lienzo, que proyectado no se lee (ver la nota de densidad
   en `irem.tex`). Así que TODOS los tamaños de aquí son los del PDF
   multiplicados por K, sin excepción, y en los elementos de despliegue el
   resultado coincide con lo que el master declara (portadilla y numerales, 44).
   Consecuencia buscada: las dos salidas parten las líneas en las mismas
   palabras, y una corrección hecha sobre una vale para la otra.

3. LAS MEDIDAS SALEN DEL MASTER O DEL PDF, NUNCA DE LA VISTA. Cada constante
   dice de dónde viene. Si cambias una, vuelve a medir con `revisar.py`.
"""
import re
import sys
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Mm, Pt

# --- Lienzo y factor de conversión -----------------------------------------
#  El lienzo del master (338.667 x 190.5 mm) contra el de Beamer (160 x 90):
#  toda medida tomada del PDF se multiplica por K para llegar aquí.
K = 338.667 / 160.0          # 2.116667
LIENZO_W, LIENZO_H = 338.667, 190.5

# --- Paleta (lámina 2 del master) ------------------------------------------
AZUL = RGBColor(0x36, 0x7C, 0xBC)
VERDE = RGBColor(0x98, 0xCE, 0x63)
GRIS_OSCURO = RGBColor(0x40, 0x40, 0x40)
GRIS_CLARO = RGBColor(0x9B, 0x9B, 0x9D)
GRIS_SUAVE = RGBColor(0xEF, 0xEE, 0xEE)
CALOR = {"bajo": RGBColor(0xF8, 0x69, 0x6B),
         "medio": RGBColor(0xFF, 0xEB, 0x84),
         "alto": RGBColor(0x63, 0xBE, 0x7B)}

FUENTE = "Montserrat"
FUENTE_MONO = "Menlo"

# --- Geometría --------------------------------------------------------------
#  Título de lámina: es el placeholder del layout «Titulo y contenido», tal
#  como lo trae el master.
TIT = (23.7, 11.0, 291.3, 15.3)

#  Cuerpo. La x y el ancho NO son los del placeholder del master (23.7 y 201.9):
#  son el área de texto del PDF convertida, porque es lo que hace que las dos
#  salidas partan las líneas igual. El placeholder del master es más angosto
#  porque sus láminas reservan la derecha para tablas y figuras.
CUERPO_X = 12.2 * K                  # 25.82
CUERPO_W = (160 - 2 * 12.2) * K      # 286.94
CUERPO_Y = 50.7                      # y del placeholder del master; en el PDF, 23.95 mm

#  Debajo de esto no va nada: es donde arranca la banda de los logotipos del
#  pie. En el PDF son los 77 mm de siempre.
AREA_FIN = 77.0 * K                  # 163.0

#  PowerPoint no pone la primera línea pegada al borde de la caja: sobre ella
#  queda el aire del interlineado. Medido contra el PDF del propio formato con
#  cuerpo de 21 pt e interlineado 1.45, la diferencia es de 4.2 mm. Se resta
#  para que el primer renglón caiga en el mismo sitio en las dos salidas. Si
#  cambias el interlineado del cuerpo, vuelve a medirlo.
AIRE_PRIMERA_LINEA = 4.2

#  Numerales: cifra, caja y paso entre filas, del layout «1_Titulo y contenido».
NUM_CIFRA = (22.3, 56.6, 26.0, 15.9)
NUM_CAJA = (48.4, 48.0, 152.7, 33.1)
NUM_PASO = 36.4
#  Recuadro de realce, del layout «2_Titulo y contenido».
REALCE = (239.2, 73.0, 77.2, 47.7)
#  Panel lateral derecho, del placeholder de cuerpo del layout «1_Titulo y contenido».
PANEL = (219.1, 49.1, 102.5, 51.7)
#  Cierre, del layout «Cierre». Solo se usa cuando el saludo no es el de fábrica.
CIERRE_SALUDO = (23.3, 37.7, 292.1, 32.4)
CIERRE_BANDA = (36.2, 80.3, 267.6, 58.7)
CIERRE_BAJADA = (84.7, 167.1, 171.2, 7.7)
#  Logotipos del pie, de cualquier layout de contenido.
LOGO_MESO = (4.2, 164.1, 44.4, 11.9)
LOGO_BID = (310.9, 166.2, 20.8, 8.4)
#  Portada: posiciones de los placeholders del layout TAPA. La autoría es la
#  excepción: el layout la pone en y=127.4 pero el master la baja a 142.2 en su
#  lámina, y esa es la que replica el PDF.
TAPA_AUTOR_Y = 142.2

#  Dónde se centra una lámina de una sola afirmación. En el PDF el bloque va
#  entre dos \vspace{\fill}, pero el formato compone con la opción `t`, así que
#  los dos rellenos no estiran igual y el bloque no queda en el medio del área:
#  medido en el PDF, el centro cae en 48 mm en una lámina con título y en 33 mm
#  en una `.plain`, que no lo tiene. Convertidos:
CENTRO_DISPLAY = 48.0 * K            # 101.6
CENTRO_DISPLAY_PLAIN = 33.0 * K      # 69.85

#  Cajas de concepto: 66 x 20 mm del PDF, dos por fila, con el aire de 4 mm.
CONC_W, CONC_H = 66 * K, 20 * K
CONC_AIRE = 4 * K
CONC_FILETE = 1.4 * K

# --- Escala tipográfica -----------------------------------------------------
#  PDF x K = el punto que declara `irem.tex`, multiplicado por el factor de
#  lienzo. Donde dice «master», ese número coincide con el que el master aplica
#  en su propio archivo, que es la comprobación de que la conversión es buena.
PT_TITULO = 34          # PDF 16 x K
PT_CUERPO = 21          # PDF 10 x K
PT_CUERPO2 = 18         # PDF 8.5 x K
PT_TABLA = 18           # PDF 8.5 x K
PT_NOTAPIE = 15         # PDF 7 x K
PT_PORTADILLA = 44      # master
PT_TAPA_TITULO = 42     # PDF 20 x K
PT_TAPA_VOLANTA = 15    # PDF 7 x K
PT_TAPA_AUTOR = 19      # PDF 9 x K
PT_TAPA_FECHA = 17      # PDF 8 x K
PT_NUM_CIFRA = 44       # master
PT_NUM_CAJA = 19        # PDF 9 x K
PT_REALCE = 19          # PDF 9 x K
PT_PANEL = 19           # PDF 9 x K
PT_CONC_T = 23          # PDF 11 x K
PT_CONC_F = 18          # PDF 8.5 x K
PT_IDEA = 42            # PDF 20 x K
PT_CIFRA = 93           # PDF 44 x K
PT_CIFRA_GLOSA = 25     # PDF 12 x K
PT_PREG_ET = 19         # PDF 9 x K
PT_PREG = 38            # PDF 18 x K
PT_GRACIAS = 55         # PDF 26 x K
PT_BAJADA = 17          # PDF 8 x K

#  Interlineados. OJO con la unidad: LaTeX declara el interlineado en puntos
#  absolutos (`\fontsize{10}{14.5}` = 14.5 pt entre líneas base) y PowerPoint lo
#  declara como múltiplo de la ALTURA NATURAL DE LÍNEA de la fuente, que en
#  Montserrat es 1.219 em, no del cuerpo. Tomar el 1.45 de LaTeX y ponerlo tal
#  cual deja las líneas un 22 % más separadas, y un punteo de cuatro renglones
#  termina 5 mm por debajo del que compone el PDF. La conversión es
#  ls = interlineado_latex / (ALTURA_LINEA x cuerpo_latex).
ALTURA_LINEA = 1.219
LS_CUERPO = 14.5 / (ALTURA_LINEA * 10)    # 1.19
LS_LISTA = 15.0 / (ALTURA_LINEA * 10)     # 1.23
LS_TAPA = 25.0 / (ALTURA_LINEA * 20)      # 1.03
LS_NUM = 11.5 / (ALTURA_LINEA * 9)        # 1.05
LS_REALCE = 12.0 / (ALTURA_LINEA * 9)     # 1.09
LS_PANEL = 12.5 / (ALTURA_LINEA * 9)      # 1.14
LS_CONC = 11.0 / (ALTURA_LINEA * 8.5)     # 1.06
#  El título va con el 0.9 del master y no con el 1.02 que daría la conversión:
#  medido, con 0.9 la primera línea cae a 0.6 mm de donde la pone el PDF, y con
#  1.02 se va al doble. Un título de dos líneas queda algo más apretado que en
#  el PDF, que es justo lo que hace el master.
LS_TITULO = 0.90
#  Aire entre puntos: 5.5 pt del PDF, convertidos.
SP_PUNTOS = Pt(5.5 * K)

MESES = ("enero", "febrero", "marzo", "abril", "mayo", "junio", "julio",
         "agosto", "septiembre", "octubre", "noviembre", "diciembre")


# ==========================================================================
#  Métricas: hacen falta para tres cosas y ninguna es decorativa: encoger la
#  volanta si no cabe en la barra, saber dónde termina un bloque para poner el
#  siguiente, y dar los anchos naturales de las columnas de una tabla.
# ==========================================================================
class Metricas:
    """Mide texto en Montserrat. Si no encuentra la fuente, estima."""

    RUTAS = (
        Path.home() / "Library/Fonts",
        Path("/Library/Fonts"),
        Path("/System/Library/Fonts/Supplemental"),
    )

    def __init__(self):
        self.fuentes = {}
        self.hay_fuente = False
        try:
            from PIL import ImageFont
        except ImportError:
            return
        for negrita in (False, True):
            nombre = "Montserrat-Bold.otf" if negrita else "Montserrat-Regular.otf"
            ruta = self._buscar(nombre)
            if ruta:
                # 1000 pt para medir una vez y escalar; así no se recarga la
                # fuente en cada medición.
                self.fuentes[negrita] = ImageFont.truetype(str(ruta), 1000)
        self.hay_fuente = len(self.fuentes) == 2

    def _buscar(self, nombre):
        for base in self.RUTAS:
            if (base / nombre).exists():
                return base / nombre
        # TeX Live la trae, y el camino en PDF ya la exige.
        import subprocess
        r = subprocess.run(["kpsewhich", nombre], capture_output=True, text=True)
        ruta = Path(r.stdout.strip()) if r.stdout.strip() else None
        return ruta if ruta and ruta.exists() else None

    def ancho_mm(self, texto, pt, negrita=False):
        """Ancho de una línea, en mm."""
        if self.hay_fuente:
            f = self.fuentes[negrita]
            return f.getlength(texto) / 1000.0 * pt * 25.4 / 72.0
        # Sin la fuente: 0.55 em por carácter, medido sobre Montserrat con
        # texto en español. Sirve para no romperse, no para calibrar.
        return len(texto) * 0.55 * pt * 25.4 / 72.0

    def lineas(self, texto, pt, ancho_mm, negrita=False):
        """Cuántos renglones ocupa el texto envuelto a ese ancho."""
        n = 0
        for parrafo in texto.split("\n"):
            palabras = parrafo.split()
            if not palabras:
                n += 1
                continue
            linea, cuenta = "", 1
            for p in palabras:
                prueba = f"{linea} {p}".strip()
                if self.ancho_mm(prueba, pt, negrita) > ancho_mm and linea:
                    cuenta += 1
                    linea = p
                else:
                    linea = prueba
            n += cuenta
        return n

    def alto_mm(self, texto, pt, ancho_mm, ls=1.0, negrita=False):
        """Alto del bloque de texto, en mm. PowerPoint compone la línea en
        1.219 em y le aplica el interlineado."""
        return self.lineas(texto, pt, ancho_mm, negrita) * pt * 1.219 * ls * 25.4 / 72.0


MET = Metricas()


# ==========================================================================
#  Lectura del .qmd
# ==========================================================================
def leer_meta(texto):
    """Encabezado YAML. No se usa un lector de YAML para no arrastrar una
    dependencia por cuatro campos de una línea."""
    m = re.match(r"^---\n(.*?)\n---\n", texto, re.S)
    if not m:
        return {}, texto
    meta = {}
    for linea in m.group(1).split("\n"):
        c = re.match(r"^([a-z-]+):\s*(.*)$", linea)
        if c:
            meta[c.group(1)] = c.group(2).strip().strip('"').strip("'")
    return meta, texto[m.end():]


def fecha_larga(valor):
    """`today`, una fecha ISO o lo que haya escrito el usuario."""
    if not valor:
        return ""
    if valor == "today":
        d = date.today()
    else:
        m = re.match(r"^(\d{4})-(\d{2})-(\d{2})$", valor)
        if not m:
            return valor
        d = date(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    return f"{d.day} de {MESES[d.month - 1]} de {d.year}"


def trocear(cuerpo):
    """Trocea el .qmd en láminas, igual que `revisar.py`: cada trozo termina
    donde empieza la lámina o la sección siguiente."""
    partes = re.split(r"^(#{1,2}) ", cuerpo, flags=re.M)[1:]
    for i in range(0, len(partes), 2):
        nivel, resto = len(partes[i]), partes[i + 1]
        titulo, _, cont = resto.partition("\n")
        yield nivel, titulo.strip(), cont


# --- Texto en línea ---------------------------------------------------------
#  **negrita**, *cursiva*, `código`. Se devuelve una lista de tramos para
#  poder darle a cada uno su formato, porque en PowerPoint el formato vive en
#  el renglón (`run`), no en el párrafo.
def tramos(texto):
    salida, resto = [], texto
    patron = re.compile(r"(\*\*(.+?)\*\*|(?<!\*)\*([^*]+?)\*(?!\*)|`([^`]+?)`)", re.S)
    pos = 0
    for m in patron.finditer(resto):
        if m.start() > pos:
            salida.append((resto[pos:m.start()], {}))
        if m.group(2) is not None:
            salida.append((m.group(2), {"negrita": True}))
        elif m.group(3) is not None:
            salida.append((m.group(3), {"cursiva": True}))
        else:
            salida.append((m.group(4), {"mono": True}))
        pos = m.end()
    if pos < len(resto):
        salida.append((resto[pos:], {}))
    return [(t, f) for t, f in salida if t] or [("", {})]


def limpiar_latex(s):
    """Lo que un texto de LaTeX trae y en PowerPoint no va."""
    s = s.replace("\\%", "%").replace("\\&", "&").replace("\\_", "_")
    s = s.replace("~", " ").replace("\\\\", "\n")
    s = re.sub(r"\\(?:textbf|textit|emph)\{([^{}]*)\}", r"\1", s)
    return s.strip()


def argumentos(cadena, pos):
    """Lee los argumentos {..} de un comando de LaTeX contando llaves, que es
    lo único que aguanta un `\\celdaCalor{alto}{50\\%}` dentro de una celda."""
    args = []
    i = pos
    while i < len(cadena) and cadena[i] == "{":
        prof, j = 0, i
        while j < len(cadena):
            if cadena[j] == "{":
                prof += 1
            elif cadena[j] == "}":
                prof -= 1
                if prof == 0:
                    break
            j += 1
        args.append(cadena[i + 1:j])
        i = j + 1
    return args, i


def parsear_tabla_latex(bloque):
    """`\\begin{tablaIrem}{m{58mm}rr} ... \\end{tablaIrem}`"""
    m = re.search(r"\\begin\{tablaIrem\}", bloque)
    if not m:
        return None
    spec, pos = argumentos(bloque, m.end())
    fin = bloque.find("\\end{tablaIrem}", pos)
    cuerpo = bloque[pos:fin]
    filas = []
    for linea in re.split(r"\\\\", cuerpo):
        if not linea.strip():
            continue
        celdas = []
        for celda in linea.split("&"):
            celda = celda.strip()
            calor = None
            c = re.search(r"\\celdaCalor", celda)
            if c:
                args, _ = argumentos(celda, c.end())
                if len(args) == 2:
                    calor, celda = args[0], args[1]
            ch = re.search(r"\\ch", celda)
            if ch:
                args, _ = argumentos(celda, ch.end())
                celda = args[0] if args else celda
            celdas.append((limpiar_latex(celda), calor))
        filas.append(celdas)
    return {"spec": spec[0] if spec else "l", "filas": filas}


def columnas_de_spec(spec):
    """`m{58mm}rr` -> [('l', 58.0), ('r', None), ('r', None)]"""
    cols, i = [], 0
    while i < len(spec):
        c = spec[i]
        if c in "mpb" and i + 1 < len(spec) and spec[i + 1] == "{":
            args, i = argumentos(spec, i + 1)
            mm = re.match(r"([\d.]+)\s*mm", args[0].strip()) if args else None
            cols.append(("l", float(mm.group(1)) if mm else None))
            continue
        if c in "lrc":
            cols.append((c, None))
        i += 1
    return cols


def parsear_tabla_markdown(lineas):
    """Tabla de Markdown: cabecera, separador con guiones y filas."""
    filas = []
    alin = []
    for i, linea in enumerate(lineas):
        celdas = [c.strip() for c in linea.strip().strip("|").split("|")]
        if i == 1 and all(re.fullmatch(r":?-{2,}:?", c) for c in celdas):
            alin = ["c" if c.startswith(":") and c.endswith(":")
                    else "r" if c.endswith(":") else "l" for c in celdas]
            continue
        filas.append([(c, None) for c in celdas])
    cols = [(a, None) for a in alin] or [("l", None)] * len(filas[0])
    return {"spec": None, "filas": filas, "cols": cols}


def parsear_lamina(cont):
    """Devuelve (bloques, notas). Los bloques van en el orden en que aparecen;
    cada uno se dibuja donde le toca."""
    bloques, notas = [], []

    # Notas del presentador: las dos formas.
    for m in re.finditer(r"::: \{\.notes\}(.*?):::", cont, re.S):
        notas.append(m.group(1).strip())
    cont = re.sub(r"::: \{\.notes\}.*?:::", "", cont, flags=re.S)

    # Bloques de LaTeX crudo. Se sacan del flujo y se parsean aparte.
    def tomar_latex(m):
        bruto = m.group(1)
        for n in re.finditer(r"\\note\{", bruto):
            args, _ = argumentos(bruto, n.end() - 1)
            if args:
                notas.append(limpiar_latex(args[0]))
        bloques.extend(parsear_latex(bruto))
        return "\n"

    cont = re.sub(r"```\{=latex\}(.*?)```", tomar_latex, cont, flags=re.S)
    # Bloques de código de otros lenguajes: en PowerPoint no se ejecutan.
    cont = re.sub(r"```.*?```", "", cont, flags=re.S)

    bloques.extend(parsear_markdown(cont))
    return bloques, "\n\n".join(n for n in notas if n)


def parsear_latex(bruto):
    """Los comandos del formato. El orden importa: `numerados` primero para que
    el contador de filas arranque en cero, como en el PDF."""
    bloques = []
    tabla = parsear_tabla_latex(bruto)
    if tabla:
        tabla["cols"] = columnas_de_spec(tabla["spec"])
        bloques.append(("tabla", tabla))

    nums = []
    for m in re.finditer(r"\\numerado", bruto):
        args, _ = argumentos(bruto, m.end())
        if len(args) == 2:
            nums.append((limpiar_latex(args[0]), limpiar_latex(args[1])))
    if nums:
        bloques.append(("numerados", nums))

    concs = []
    for m in re.finditer(r"\\concepto", bruto):
        args, _ = argumentos(bruto, m.end())
        if len(args) == 2:
            concs.append((limpiar_latex(args[0]), limpiar_latex(args[1])))
    if concs:
        bloques.append(("conceptos", concs))

    for cmd, clave in (("realce", "realce"), ("panelDerecho", "panel"),
                       ("ideaGrande", "idea"), ("pregunta", "pregunta"),
                       ("notaPie", "notapie"), ("cintilla", "cintilla")):
        for m in re.finditer(rf"\\{cmd}(?![a-zA-Z])", bruto):
            args, _ = argumentos(bruto, m.end())
            if args:
                bloques.append((clave, limpiar_latex(args[0])))

    for m in re.finditer(r"\\cifra(?![a-zA-Z])", bruto):
        args, _ = argumentos(bruto, m.end())
        if len(args) == 2:
            bloques.append(("cifra", (limpiar_latex(args[0]), limpiar_latex(args[1]))))

    for m in re.finditer(r"\\laminaGracias(\[([^\]]*)\])?", bruto):
        bloques.append(("gracias", m.group(2)))

    if re.search(r"\\logosPie", bruto):
        bloques.append(("logospie", None))

    return bloques


def parsear_markdown(cont):
    """Punteos, párrafos, tablas, imágenes y columnas."""
    bloques = []
    lineas = cont.split("\n")
    i = 0
    while i < len(lineas):
        linea = lineas[i]
        s = linea.strip()

        if not s:
            i += 1
            continue

        # Columnas
        if re.match(r"^:{2,}\s*\{\.columns\}", s):
            prof, j, dentro = 1, i + 1, []
            while j < len(lineas) and prof > 0:
                t = lineas[j].strip()
                if re.match(r"^:{2,}\s*\{", t):
                    prof += 1
                elif re.match(r"^:{3,}$|^:{2,}$", t):
                    prof -= 1
                    if prof == 0:
                        break
                dentro.append(lineas[j])
                j += 1
            cols = []
            actual, ancho = None, None
            for t in dentro:
                m = re.match(r"^:{2,}\s*\{\.column\s+width=\"(\d+)%\"\}", t.strip())
                if m:
                    if actual is not None:
                        cols.append((ancho, "\n".join(actual)))
                    actual, ancho = [], int(m.group(1)) / 100.0
                elif re.match(r"^:{2,}$", t.strip()):
                    if actual is not None:
                        cols.append((ancho, "\n".join(actual)))
                        actual, ancho = None, None
                elif actual is not None:
                    actual.append(t)
            if actual is not None:
                cols.append((ancho, "\n".join(actual)))
            bloques.append(("columnas", cols))
            i = j + 1
            continue

        # Marcas de bloque sueltas
        if re.match(r"^:{2,}", s):
            i += 1
            continue

        # Imagen
        m = re.match(r"^!\[([^\]]*)\]\(([^)]+)\)", s)
        if m:
            bloques.append(("imagen", m.group(2)))
            i += 1
            continue

        # Tabla de Markdown
        if s.startswith("|"):
            j = i
            filas = []
            while j < len(lineas) and lineas[j].strip().startswith("|"):
                filas.append(lineas[j])
                j += 1
            if len(filas) >= 2:
                bloques.append(("tabla", parsear_tabla_markdown(filas)))
                i = j
                continue

        # Punteos y listas numeradas
        if re.match(r"^\s*([-*]|\d+\.)\s+", linea):
            items = []
            while i < len(lineas):
                l2 = lineas[i]
                m2 = re.match(r"^(\s*)([-*]|\d+\.)\s+(.*)$", l2)
                if m2:
                    nivel = 1 if len(m2.group(1)) >= 2 else 0
                    numerada = not m2.group(2) in ("-", "*")
                    items.append([nivel, numerada, m2.group(3).strip()])
                elif l2.strip() and items:
                    items[-1][2] += " " + l2.strip()   # punto que sigue en el renglón siguiente
                else:
                    break
                i += 1
            bloques.append(("punteos", items))
            continue

        # Párrafo
        parrafo = [s]
        i += 1
        while i < len(lineas) and lineas[i].strip() and not re.match(
                r"^\s*([-*]|\d+\.)\s+|^\||^!\[|^:{2,}", lineas[i]):
            parrafo.append(lineas[i].strip())
            i += 1
        bloques.append(("parrafo", " ".join(parrafo)))

    return bloques


# ==========================================================================
#  Escritura del .pptx
# ==========================================================================
def pPr(p):
    return p._p.get_or_add_pPr()


def sin_vineta(p, sangria=0.0):
    """El primer nivel de este master no lleva viñeta ni sangría: va al ras
    del título. Lo declara el layout con `buNone` y lo confirma el PDF."""
    e = pPr(p)
    e.set("marL", str(int(sangria * 36000)))
    e.set("indent", "0")
    e.append(e.makeelement(qn("a:buNone"), {}))


def con_vineta(p, sangria, color=GRIS_OSCURO):
    """El punto aparece recién en el segundo nivel, y en gris oscuro."""
    e = pPr(p)
    e.set("marL", str(int(sangria * 36000)))
    e.set("indent", str(-int(3.0 * K * 36000)))
    clr = e.makeelement(qn("a:buClr"), {})
    srgb = clr.makeelement(qn("a:srgbClr"), {"val": str(color)})
    clr.append(srgb)
    e.append(clr)
    fnt = e.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    e.append(fnt)
    e.append(e.makeelement(qn("a:buChar"), {"char": "\u2022"}))


def con_numero(p, sangria):
    e = pPr(p)
    e.set("marL", str(int(sangria * 36000)))
    e.set("indent", str(-int(sangria * 36000)))
    e.append(e.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"}))


def caja(lamina, x, y, w, h, ajuste=True):
    """Caja de texto sin márgenes internos. Sin `word_wrap` explícito el texto
    se compone en una sola línea y se sale de la lámina por la derecha."""
    tb = lamina.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = tb.text_frame
    tf.word_wrap = ajuste
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb


def escribir(p, contenido, pt, color=GRIS_OSCURO, negrita=False, cursiva=False,
             fuente=FUENTE):
    """Escribe un párrafo, tramo por tramo. Aquí es donde se fija a mano lo
    que los layouts declaran mal."""
    for texto, f in (contenido if isinstance(contenido, list) else tramos(contenido)):
        r = p.add_run()
        r.text = texto
        r.font.size = Pt(pt)
        r.font.bold = negrita or f.get("negrita", False)
        r.font.italic = cursiva or f.get("cursiva", False)
        r.font.name = FUENTE_MONO if f.get("mono") else fuente
        r.font.color.rgb = color


def texto_simple(lamina, x, y, w, h, contenido, pt, **kw):
    alineacion = kw.pop("alineacion", None)
    ancla = kw.pop("ancla", None)
    ls = kw.pop("ls", None)
    tb = caja(lamina, x, y, w, h)
    p = tb.text_frame.paragraphs[0]
    if alineacion:
        p.alignment = alineacion
    if ls:
        p.line_spacing = ls
    if ancla:
        tb.text_frame.vertical_anchor = ancla
    escribir(p, contenido, pt, **kw)
    return tb


def quitar_placeholders(lamina, salvo=()):
    """Un placeholder vacío no se proyecta, pero en edición muestra su «haga
    clic para agregar texto» y estorba. Los que no se usan se van."""
    for ph in list(lamina.placeholders):
        if ph.placeholder_format.idx not in salvo:
            e = ph._element
            e.getparent().remove(e)


def sin_graficos_heredados(lamina):
    """Lámina `.plain`: el fondo y la franja se quedan (vienen del `p:bg` del
    slide master) pero los logotipos del pie, que son formas del layout, no.
    Es lo mismo que «Ocultar gráficos de fondo» en el menú de PowerPoint."""
    lamina._element.set("showMasterSp", "0")


# --- Tablas ------------------------------------------------------------------
SIN_ESTILO = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"   # «Sin estilo, sin cuadrícula»


def estilo_tabla(tabla):
    """PowerPoint le pone a toda tabla nueva un estilo azul con bandas. Hay que
    apagarlo entero antes de pintar la cabecera del formato, o queda una tabla
    con dos estilos encima."""
    tp = tabla._tbl.find(qn("a:tblPr"))
    if tp is None:
        return
    for atr in ("bandRow", "bandCol", "firstCol", "lastRow", "lastCol"):
        tp.set(atr, "0")
    tp.set("firstRow", "0")
    for hijo in list(tp):
        if hijo.tag == qn("a:tableStyleId"):
            tp.remove(hijo)
    e = tp.makeelement(qn("a:tableStyleId"), {})
    e.text = SIN_ESTILO
    tp.append(e)


def filete(celda, lado="B", color=GRIS_CLARO, grosor=0.15):
    """Filete horizontal fino y gris, y ninguno vertical. Va antes que el
    relleno porque en `a:tcPr` los bordes van primero; al revés, PowerPoint
    descarta el archivo."""
    tc = celda._tc.get_or_add_tcPr()
    ln = tc.makeelement(qn(f"a:ln{lado}"), {"w": str(int(grosor * 12700)), "cap": "flat",
                                            "cmpd": "sng", "algn": "ctr"})
    fill = ln.makeelement(qn("a:solidFill"), {})
    srgb = fill.makeelement(qn("a:srgbClr"), {"val": str(color)})
    fill.append(srgb)
    ln.append(fill)
    tc.insert(0, ln)


def poner_tabla(lamina, datos, y):
    """Tabla del formato: cabecera azul con texto gris suave en negrita, cuerpo
    sin relleno. El ancho de cada columna sale del contenido, como en el PDF,
    para que la tabla no se estire a todo lo ancho cuando no hace falta."""
    filas = datos["filas"]
    cols = datos.get("cols") or [("l", None)] * len(filas[0])
    n_col = max(len(f) for f in filas)
    while len(cols) < n_col:
        cols.append(("l", None))

    # Filetes: `tablaIrem` no lleva ninguno (su separación es la banda azul y el
    # aire), y una tabla de Markdown lleva los tres de booktabs, que es como la
    # compone pandoc en el PDF. Se distingue por el `spec`, que solo trae la de
    # LaTeX.
    booktabs = datos.get("spec") is None

    relleno = 2.4 * K * 2      # el tabcolsep del PDF, a los dos lados
    anchos = []
    for c in range(n_col):
        fijo = cols[c][1]
        if fijo:
            anchos.append(fijo * K)
            continue
        ancho = 0.0
        for fi, fila in enumerate(filas):
            if c < len(fila):
                ancho = max(ancho, MET.ancho_mm(fila[c][0], PT_TABLA, negrita=(fi == 0)))
        # Un 2 % de holgura más un milímetro: PowerPoint decide el corte de
        # línea con un poco más de aire que la medición, y sin esto una
        # cabecera que cabía justo se parte en dos renglones.
        anchos.append(ancho * 1.02 + relleno + 1.0)
    total = sum(anchos)
    if total > CUERPO_W:                      # no cabe: se encoge proporcional
        anchos = [a * CUERPO_W / total for a in anchos]
        total = CUERPO_W
    x = CUERPO_X + (CUERPO_W - total) / 2     # centrada, como el \begin{center}

    alto_fila = PT_TABLA * 1.2 * 1.45 * 25.4 / 72.0    # arraystretch del PDF
    forma = lamina.shapes.add_table(len(filas), n_col, Mm(x), Mm(y),
                                    Mm(total), Mm(alto_fila * len(filas)))
    tabla = forma.table
    estilo_tabla(tabla)
    for c, a in enumerate(anchos):
        tabla.columns[c].width = Mm(a)
    alto_total = 0.0
    for fi, fila in enumerate(filas):
        lineas = 1
        for c in range(n_col):
            if c < len(fila):
                lineas = max(lineas, MET.lineas(fila[c][0], PT_TABLA,
                                                anchos[c] - relleno, fi == 0))
        alto = alto_fila * lineas
        tabla.rows[fi].height = Mm(alto)
        alto_total += alto
        for c in range(n_col):
            celda = tabla.cell(fi, c)
            celda.margin_left = celda.margin_right = Mm(2.4 * K)
            celda.margin_top = celda.margin_bottom = 0
            celda.vertical_anchor = MSO_ANCHOR.MIDDLE
            if booktabs:
                if fi == 0:
                    filete(celda, "T")
                if fi in (0, len(filas) - 1):
                    filete(celda, "B")
            if fi == 0:
                celda.fill.solid()
                celda.fill.fore_color.rgb = AZUL
            else:
                celda.fill.background()
            texto, calor = fila[c] if c < len(fila) else ("", None)
            if calor and calor in CALOR:
                celda.fill.solid()
                celda.fill.fore_color.rgb = CALOR[calor]
            p = celda.text_frame.paragraphs[0]
            alin = cols[c][0]
            p.alignment = {"r": PP_ALIGN.RIGHT, "c": PP_ALIGN.CENTER}.get(alin, PP_ALIGN.LEFT)
            escribir(p, texto, PT_TABLA,
                     color=GRIS_SUAVE if fi == 0 else GRIS_OSCURO,
                     negrita=(fi == 0))
    return y + alto_total


# --- Elementos del formato ---------------------------------------------------
def poner_numerados(lamina, nums):
    """Cifra verde y caja verde, en posición absoluta: la fila n cae donde el
    master la pone. Máximo tres; el master nunca pone una cuarta y no cabe."""
    for i, (cifra, texto) in enumerate(nums[:3]):
        y = NUM_CAJA[1] + i * NUM_PASO
        forma = lamina.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Mm(NUM_CAJA[0]),
                                        Mm(y), Mm(NUM_CAJA[2]), Mm(NUM_CAJA[3]))
        forma.fill.solid()
        forma.fill.fore_color.rgb = VERDE
        forma.line.fill.background()
        forma.shadow.inherit = False
        forma.adjustments[0] = 0.04
        tf = forma.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Mm(3 * K)
        tf.margin_top = tf.margin_bottom = Mm(1.5 * K)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = LS_NUM
        escribir(p, texto, PT_NUM_CAJA, color=GRIS_OSCURO)
        texto_simple(lamina, NUM_CIFRA[0], y + (NUM_CAJA[3] - NUM_CIFRA[3]) / 2,
                     NUM_CIFRA[2], NUM_CIFRA[3], cifra, PT_NUM_CIFRA,
                     color=VERDE, negrita=True, ancla=MSO_ANCHOR.MIDDLE)


def poner_realce(lamina, texto):
    """Recuadro azul de contorno a la derecha de los numerales."""
    forma = lamina.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Mm(REALCE[0]),
                                    Mm(REALCE[1]), Mm(REALCE[2]), Mm(REALCE[3]))
    forma.fill.background()
    forma.line.color.rgb = AZUL
    forma.line.width = Pt(0.35 * K * 72 / 25.4)
    forma.shadow.inherit = False
    forma.adjustments[0] = 0.03
    tf = forma.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Mm(2 * K)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = LS_REALCE
    escribir(p, texto, PT_REALCE, color=AZUL)


def poner_conceptos(lamina, concs, y):
    """Rejilla de cajas: filete verde a la izquierda, fondo gris suave, título
    en azul. Dos por fila, que es lo que cabe en el ancho de texto."""
    for i, (termino, frase) in enumerate(concs[:4]):
        fila, col = divmod(i, 2)
        x = CUERPO_X if col == 0 else CUERPO_X + CUERPO_W - CONC_W
        yy = y + fila * (CONC_H + CONC_AIRE)
        fondo = lamina.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Mm(x), Mm(yy),
                                        Mm(CONC_W), Mm(CONC_H))
        fondo.fill.solid()
        fondo.fill.fore_color.rgb = GRIS_SUAVE
        fondo.line.fill.background()
        fondo.shadow.inherit = False
        fondo.adjustments[0] = 0.04
        fondo.text_frame.word_wrap = True
        # El filete de un solo lado no existe en PowerPoint: va como una forma
        # aparte, que además queda editable.
        filete = lamina.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(x), Mm(yy),
                                         Mm(CONC_FILETE), Mm(CONC_H))
        filete.fill.solid()
        filete.fill.fore_color.rgb = VERDE
        filete.line.fill.background()
        filete.shadow.inherit = False
        tf = fondo.text_frame
        tf.margin_left = Mm(3 * K)
        tf.margin_right = Mm(2.5 * K)
        tf.margin_top = Mm(2 * K)
        tf.margin_bottom = Mm(1 * K)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        escribir(p, termino, PT_CONC_T, color=AZUL, negrita=True)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.line_spacing = LS_CONC
        p2.space_before = Pt(1.2 * K)
        escribir(p2, frase, PT_CONC_F, color=GRIS_OSCURO)
    filas = (min(len(concs), 4) + 1) // 2
    return y + filas * (CONC_H + CONC_AIRE)


def poner_logos_pie(lamina, logos):
    for ruta, (x, y, w, h) in ((logos / "logo-mesoamerica.png", LOGO_MESO),
                               (logos / "logo-bid.png", LOGO_BID)):
        if ruta.exists():
            f = lamina.shapes.add_picture(str(ruta), Mm(x), Mm(y), Mm(w), Mm(h))
            # El nombre no es cosmético: `revisar.py` se salta por él las piezas
            # que sí van dentro de la banda del pie, y de paso quien abra el
            # archivo las encuentra por nombre en el panel de selección.
            f.name = "pie-" + ruta.stem


def poner_cintilla(lamina, nombre, logos):
    """`\\cintilla{comite}` o `\\cintilla{socios}`: al pie, centrada, 120 mm
    del PDF a 7 mm del borde."""
    ruta = logos / f"cintilla-{nombre}-esp.png"
    if not ruta.exists():
        return
    from PIL import Image
    w = 120 * K
    with Image.open(ruta) as im:
        h = w * im.size[1] / im.size[0]
    f = lamina.shapes.add_picture(str(ruta), Mm((LIENZO_W - w) / 2),
                                  Mm(LIENZO_H - 7 * K - h), Mm(w), Mm(h))
    f.name = "pie-cintilla"


def poner_gracias(lamina, texto, logos):
    """Cierre: saludo, banda de logotipos y bajada.

    Se dibuja aunque el master traiga el layout «Cierre» ya montado, por dos
    razones. Una, que el saludo del layout NO se puede editar desde la lámina,
    y quien recibe un .pptx espera poder cambiarlo. Y dos, que la bajada del
    layout va en 12 pt, más pequeña que la del PDF."""
    sin_graficos_heredados(lamina)
    saludo = texto_simple(lamina, *CIERRE_SALUDO, texto, PT_GRACIAS, color=AZUL,
                          negrita=True, alineacion=PP_ALIGN.CENTER)
    saludo.name = "cierre-saludo"
    banda = logos / "cierre-logos.png"
    if banda.exists():
        f = lamina.shapes.add_picture(str(banda), Mm(CIERRE_BANDA[0]), Mm(CIERRE_BANDA[1]),
                                      Mm(CIERRE_BANDA[2]), Mm(CIERRE_BANDA[3]))
        f.name = "cierre-banda"
    bajada = texto_simple(lamina, *CIERRE_BAJADA,
                          "Iniciativa Regional para la Eliminación de la Malaria",
                          PT_BAJADA, color=GRIS_OSCURO, alineacion=PP_ALIGN.CENTER)
    bajada.name = "cierre-bajada"


def poner_punteos(lamina, items, y, ancho=None, pt=None):
    """Punteos y listas numeradas. Primer nivel al ras y sin viñeta, segundo
    con punto gris; lo que separa un punto del siguiente es el aire."""
    ancho = ancho or CUERPO_W
    pt = pt or PT_CUERPO
    alto = 0.0
    for nivel, _, texto in items:
        alto += MET.alto_mm(texto, pt if nivel == 0 else PT_CUERPO2,
                            ancho - (6.0 * K if nivel else 0), LS_LISTA)
    alto += SP_PUNTOS.pt * (len(items) - 1) * 25.4 / 72.0
    tb = caja(lamina, CUERPO_X, y - AIRE_PRIMERA_LINEA, ancho, max(alto + 4, 10))
    tf = tb.text_frame
    for i, (nivel, numerada, texto) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = LS_LISTA
        if i:
            p.space_before = SP_PUNTOS
        if numerada:
            con_numero(p, 5.0 * K)
        elif nivel:
            con_vineta(p, 6.0 * K)
        else:
            sin_vineta(p)
        escribir(p, texto, pt if nivel == 0 else PT_CUERPO2)
    return y + alto


def poner_parrafo(lamina, texto, y, ancho=None, pt=None):
    ancho = ancho or CUERPO_W
    pt = pt or PT_CUERPO
    alto = MET.alto_mm(texto, pt, ancho, LS_CUERPO)
    tb = caja(lamina, CUERPO_X, y - AIRE_PRIMERA_LINEA, ancho, max(alto + 4, 8))
    p = tb.text_frame.paragraphs[0]
    p.line_spacing = LS_CUERPO
    sin_vineta(p)
    escribir(p, texto, pt)
    return y + alto


def poner_imagen(lamina, ruta, y, base):
    from PIL import Image
    archivo = (base / ruta) if not Path(ruta).is_absolute() else Path(ruta)
    if not archivo.exists():
        print(f"  aviso: no encuentro la imagen {ruta}")
        return y
    with Image.open(archivo) as im:
        prop = im.size[1] / im.size[0]
    w = CUERPO_W
    h = w * prop
    disponible = AREA_FIN - y - 2
    if h > disponible:
        h = disponible
        w = h / prop
    lamina.shapes.add_picture(str(archivo), Mm(CUERPO_X + (CUERPO_W - w) / 2),
                              Mm(y), Mm(w), Mm(h))
    return y + h


def poner_bloques(lamina, bloques, y, base, logos, ancho=None, plain=False):
    """Dibuja los bloques de una lámina y devuelve dónde quedó el cursor."""
    centro = CENTRO_DISPLAY_PLAIN if plain else CENTRO_DISPLAY
    for clase, datos in bloques:
        if clase == "punteos":
            y = poner_punteos(lamina, datos, y, ancho)
        elif clase == "parrafo":
            y = poner_parrafo(lamina, datos, y, ancho)
        elif clase == "tabla":
            y = poner_tabla(lamina, datos, y) + 1.5 * K
        elif clase == "imagen":
            y = poner_imagen(lamina, datos, y, base) + 1.5 * K
        elif clase == "numerados":
            poner_numerados(lamina, datos)
        elif clase == "realce":
            poner_realce(lamina, datos)
        elif clase == "panel":
            texto_simple(lamina, PANEL[0], PANEL[1] - AIRE_PRIMERA_LINEA, PANEL[2],
                         PANEL[3], datos, PT_PANEL, ls=LS_PANEL)
        elif clase == "conceptos":
            y = poner_conceptos(lamina, datos, y)
        elif clase == "notapie":
            # Si la lámina lleva un bloque de despliegue, el cursor no avanzó y
            # la nota tiene que irse al fondo del área útil, que es donde la
            # deja el \vfill del PDF.
            alto = MET.alto_mm(datos, PT_NOTAPIE, CUERPO_W, 1.29)
            yy = y if y > CUERPO_Y else AREA_FIN - alto - 2
            texto_simple(lamina, CUERPO_X, yy, CUERPO_W, alto + 4, datos,
                         PT_NOTAPIE, color=GRIS_CLARO)
            y = yy + alto
        elif clase == "idea":
            alto = MET.alto_mm(datos, PT_IDEA, CUERPO_W, 1.07, negrita=True)
            texto_simple(lamina, CUERPO_X, centro - alto / 2, CUERPO_W, alto + 4,
                         datos, PT_IDEA, color=AZUL, negrita=True,
                         alineacion=PP_ALIGN.CENTER, ls=1.07)
        elif clase == "cifra":
            dato, glosa = datos
            alto_d = MET.alto_mm(dato, PT_CIFRA, CUERPO_W, 1.0, negrita=True)
            alto_g = MET.alto_mm(glosa, PT_CIFRA_GLOSA, CUERPO_W, 1.0)
            aire = 0.35 * PT_CIFRA * 25.4 / 72.0
            arriba = centro - (alto_d + aire + alto_g) / 2
            texto_simple(lamina, CUERPO_X, arriba, CUERPO_W, alto_d + 4, dato,
                         PT_CIFRA, color=VERDE, negrita=True, alineacion=PP_ALIGN.CENTER)
            texto_simple(lamina, CUERPO_X, arriba + alto_d + aire, CUERPO_W,
                         alto_g + 4, glosa, PT_CIFRA_GLOSA, color=GRIS_CLARO,
                         alineacion=PP_ALIGN.CENTER)
        elif clase == "pregunta":
            alto_e = MET.alto_mm("PREGUNTA AL EQUIPO", PT_PREG_ET, CUERPO_W, 1.0, True)
            alto_t = MET.alto_mm(datos, PT_PREG, CUERPO_W, 1.05, negrita=True)
            aire = 0.7 * PT_PREG_ET * 25.4 / 72.0
            arriba = centro - (alto_e + aire + alto_t) / 2
            texto_simple(lamina, CUERPO_X, arriba, CUERPO_W, alto_e + 3,
                         "PREGUNTA AL EQUIPO", PT_PREG_ET, color=VERDE,
                         negrita=True, alineacion=PP_ALIGN.CENTER)
            texto_simple(lamina, CUERPO_X, arriba + alto_e + aire, CUERPO_W,
                         alto_t + 4, datos, PT_PREG, color=GRIS_OSCURO,
                         negrita=True, alineacion=PP_ALIGN.CENTER, ls=1.05)
        elif clase == "gracias":
            poner_gracias(lamina, datos or "¡Muchas gracias!", logos)
        elif clase == "cintilla":
            poner_cintilla(lamina, datos, logos)
        elif clase == "logospie":
            poner_logos_pie(lamina, logos)
        elif clase == "columnas":
            for i, (frac, contenido) in enumerate(datos):
                frac = frac or 1.0 / max(len(datos), 1)
                ancho_col = CUERPO_W * frac
                x_col = CUERPO_X + sum(CUERPO_W * (c[0] or 0) for c in datos[:i])
                sub, _ = parsear_lamina(contenido)
                cy = y
                for c_clase, c_datos in sub:
                    if c_clase == "parrafo":
                        alto = MET.alto_mm(c_datos, PT_CUERPO, ancho_col, LS_CUERPO)
                        tb = caja(lamina, x_col, cy - AIRE_PRIMERA_LINEA, ancho_col,
                                  max(alto + 4, 8))
                        p = tb.text_frame.paragraphs[0]
                        p.line_spacing = LS_CUERPO
                        sin_vineta(p)
                        escribir(p, c_datos, PT_CUERPO)
                        cy += alto + 2
                    elif c_clase == "punteos":
                        cy = poner_punteos(lamina, c_datos, cy, ancho_col) + 2
    return y


def notas(lamina, texto):
    if not texto:
        return
    lamina.notes_slide.notes_text_frame.text = texto


# ==========================================================================
def main():
    if len(sys.argv) != 2:
        sys.exit("Uso: ./renderizar-pptx.py <archivo>.qmd")
    qmd = Path(sys.argv[1]).resolve()
    if not qmd.exists():
        sys.exit(f"No existe: {qmd}")
    base = qmd.parent
    plantilla = base / "plantilla-irem.pptx"
    if not plantilla.exists():
        sys.exit("Falta plantilla-irem.pptx en la carpeta del .qmd.\n"
                 "Cópiala de la skill: cp \"$BASE/plantilla/plantilla-irem.pptx\" .")
    logos = base / "_extensions" / "irem" / "logos"

    meta, cuerpo = leer_meta(qmd.read_text())
    prs = Presentation(str(plantilla))
    L = {l.name: l for l in prs.slide_layouts}

    # --- Portada ------------------------------------------------------------
    #  Todo lo demás lo trae el layout: fotografía, cintilla de donantes, barra
    #  verde, logotipo y fondo.
    tapa = prs.slides.add_slide(L["TAPA"])
    for ph in tapa.placeholders:
        idx = ph.placeholder_format.idx
        tf = ph.text_frame
        tf.word_wrap = True
        # Todo placeholder de PowerPoint trae un margen interno de 0.1 pulgada.
        # Con él, el título de portada arrancaría 2.5 mm a la derecha de donde
        # lo pone el PDF y no quedaría alineado con el logotipo de arriba.
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        if idx == 0:                                   # título
            p.line_spacing = LS_TAPA
            escribir(p, meta.get("title", ""), PT_TAPA_TITULO, color=GRIS_OSCURO,
                     negrita=True)
        elif idx == 2:                                 # volanta, sobre la barra
            p.alignment = PP_ALIGN.CENTER
            volanta = meta.get("subtitle", "")
            pt = PT_TAPA_VOLANTA
            # La barra mide lo que mide: 107 mm menos su aire. Si la volanta no
            # cabe en un renglón, se encoge; envuelta caería fuera de la barra.
            while pt > 7 and MET.ancho_mm(volanta, pt, negrita=True) > 99.5:
                pt -= 0.5
            escribir(p, volanta, pt, color=GRIS_SUAVE, negrita=True)
        else:                                          # autoría y fecha
            # Las cuatro medidas, no solo la que cambia: si a un placeholder
            # que heredaba su posición del layout le fijas una sola, las otras
            # tres se escriben en cero y PowerPoint parte el texto letra por
            # letra contra una caja sin ancho.
            # El ancho es el del PDF (60 mm), no el del placeholder: con los
            # 83.5 mm del layout, un nombre de tres palabras se parte en dos.
            ph.left, ph.top = Mm(16.27), Mm(TAPA_AUTOR_Y)
            ph.width, ph.height = Mm(60 * K), Mm(20.1)
            escribir(p, meta.get("author", ""), PT_TAPA_AUTOR, color=GRIS_CLARO)
            fecha = fecha_larga(meta.get("date", ""))
            if fecha:
                p2 = tf.add_paragraph()
                p2.space_before = Pt(0.8 * K)
                escribir(p2, fecha, PT_TAPA_FECHA, color=GRIS_CLARO)

    # --- Láminas ------------------------------------------------------------
    cuenta = {"portadillas": 0, "contenido": 0, "notas": 0}
    for nivel, titulo, cont in trocear(cuerpo):
        if nivel == 1:
            lamina = prs.slides.add_slide(L["Portadillas"])
            for ph in lamina.placeholders:
                tf = ph.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                escribir(p, titulo, PT_PORTADILLA, color=AZUL, negrita=True)
            cuenta["portadillas"] += 1
            continue

        bloques, nota = parsear_lamina(cont)
        plain = bool(re.match(r"^\{.*\.plain.*\}$", titulo)) or titulo.startswith("{")
        solo_gracias = [b for b in bloques if b[0] == "gracias" and not b[1]]

        if solo_gracias:
            # El cierre de fábrica es el layout entero: saludo, banda de
            # logotipos y bajada. No hay nada que escribir.
            lamina = prs.slides.add_slide(L["Cierre"])
            quitar_placeholders(lamina)
            notas(lamina, nota)
            cuenta["contenido"] += 1
            cuenta["notas"] += bool(nota)
            continue

        lamina = prs.slides.add_slide(L["Titulo y contenido"])
        y = CUERPO_Y
        if plain:
            quitar_placeholders(lamina)
            sin_graficos_heredados(lamina)
        else:
            quitar_placeholders(lamina, salvo=(0,))
            ph = lamina.placeholders[0]
            ph.left, ph.top, ph.width, ph.height = (Mm(TIT[0]), Mm(TIT[1]),
                                                    Mm(TIT[2]), Mm(TIT[3]))
            tf = ph.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            p = tf.paragraphs[0]
            p.line_spacing = LS_TITULO
            escribir(p, titulo, PT_TITULO, color=AZUL, negrita=True)
            # Un título de dos renglones empuja el cuerpo, igual que en el PDF.
            renglones = MET.lineas(titulo, PT_TITULO, TIT[2], negrita=True)
            if renglones > 1:
                y += (renglones - 1) * PT_TITULO * 1.2 * LS_TITULO * 25.4 / 72.0

        poner_bloques(lamina, bloques, y, base, logos, plain=plain)
        notas(lamina, nota)
        cuenta["contenido"] += 1
        cuenta["notas"] += bool(nota)

    salida = qmd.with_suffix(".pptx")
    prs.save(str(salida))
    print(f"\n{salida.name}   {len(prs.slides.__iter__.__self__._sldIdLst)} láminas "
          f"({cuenta['contenido']} de contenido, {cuenta['portadillas']} portadillas, "
          f"{cuenta['notas']} con nota)")
    if not MET.hay_fuente:
        print("  aviso: no encontré Montserrat instalada, así que los anchos van")
        print("         estimados. Instálala (Google Fonts o tlmgr install montserrat)")
        print("         y vuelve a correr esto antes de entregar.")
    print("\nAhora revísalo:  ./revisar.py " + qmd.name)


if __name__ == "__main__":
    main()
