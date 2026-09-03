#!/usr/bin/env -S uv run --with pillow --with python-pptx --script
"""
Revisa una presentación IREM antes de entregarla.

    ./revisar.py <archivo>.qmd

Comprueba tres cosas distintas y las reporta por separado:

  FORMATO   sobre los PNG de las láminas del PDF: que nada invada la banda de
            los logotipos ni se salga por los lados, y que la portada sangre.
  POWERPOINT sobre el .pptx, si lo hay: que ninguna forma se salga del área
            útil, que el texto quepa en su caja y que ningún renglón se haya
            quedado sin la tipografía o el color del formato.
  CONTENIDO sobre el .qmd: palabras por lámina, viñetas, títulos repetidos,
            marcadores sin resolver, cobertura de notas y, si es una propuesta
            a un externo, que estén las siete secciones y en orden.

No sustituye mirar las láminas: detecta lo mecánico, no si la presentación se
entiende. Lo que no puede medir, lo dice.
"""
import re
import subprocess
import sys
from pathlib import Path

TOPE_PALABRAS = 40
TOPE_VINETAS = 4
TOPE_CARACTERES_VINETA = 118  # unos dos renglones a 10 pt en 137.6 mm
MIN_LAMINAS_POR_SECCION = 2   # menos que esto y la sección no se gana su portadilla
LAMINAS_POR_DISPLAY = 6       # una lámina de una sola frase por cada seis de contenido
TOPE_VACIAS = 0.15            # proporción máxima de láminas con una frase o menos
TINTA_MINIMA = 2.0            # % de tinta bajo el cual la lámina es "una frase"

# --- Medidas del lienzo de PowerPoint --------------------------------------
K = 338.667 / 160.0           # el mismo factor del generador
AREA_FIN_PPT = 77.0 * K       # 163.0 mm: aquí empieza la banda de los logotipos
MARGEN_IZQ_PPT = 9.5 * K      # los mismos márgenes que vigila el PDF
#  El derecho no es el del cuerpo: el panel lateral del master llega a 321.6 mm
#  (su placeholder mide 102.5 de ancho desde 219.1), y es parte del formato.
MARGEN_DER_PPT = 322.5
ALTURA_LINEA = 1.219          # altura natural de línea de Montserrat, en em

DISPLAY = r"\\(?:ideaGrande|cifra|pregunta)\b"

ESPINA = [
    "La solicitud", "Lo que sabemos", "Lo que vemos", "Lo que consideramos",
    "Lo que buscamos", "Lo que proponemos", "Siguientes pasos",
]

# Comandos del formato cuyo texto SÍ se proyecta y por lo tanto cuenta.
CMD_CON_TEXTO = r"\\(?:numerado|concepto|realce|ideaGrande|pregunta|cifra|panelDerecho)"


def laminas(cuerpo):
    """Trocea el .qmd en láminas. Cada trozo termina donde empieza la lámina o
    la sección siguiente: sin ese corte, el encabezado de la sección siguiente
    se cuenta como texto de la lámina anterior e infla el conteo."""
    partes = re.split(r"^(#{1,2}) ", cuerpo, flags=re.M)[1:]
    for i in range(0, len(partes), 2):
        nivel, resto = partes[i], partes[i + 1]
        titulo, _, cont = resto.partition("\n")
        yield len(nivel), titulo.strip(), cont


def texto_visible(cont):
    """Lo que se proyecta: fuera notas, bloques de código y marcas de bloque.
    El texto que va dentro de los comandos del formato sí se proyecta."""
    en_latex = " ".join(
        m.group(1) or ""
        for m in re.finditer(CMD_CON_TEXTO + r"(?:\[[^\]]*\])?\{[^{}]*\}(?:\{([^{}]*)\})?", cont)
    )
    vis = re.sub(r"::: \{\.notes\}.*?:::", "", cont, flags=re.S)
    vis = re.sub(r"```.*?```", "", vis, flags=re.S)
    vis = re.sub(r"^:::+.*$", "", vis, flags=re.M)
    return vis, en_latex


def palabras(vis, en_latex):
    """Cuenta palabras. Quita la marca «- » de viñeta y la de énfasis, que no
    son palabras que nadie lea."""
    t = re.sub(r"^\s*[-*]\s+", " ", vis, flags=re.M)
    t = (t + " " + en_latex).replace("*", " ").replace("\\%", "%")
    return len([w for w in t.split() if any(c.isalnum() for c in w)])


def revisar_contenido(qmd):
    s = qmd.read_text()
    partes = s.split("---", 2)
    cuerpo = partes[2] if len(partes) > 2 else s
    fallos, avisos = [], []

    secciones = [t for n, t, _ in laminas(cuerpo) if n == 1]
    # Qué páginas del PDF llevan cintilla: esa sí baja a la banda del pie, por
    # diseño, y si no se avisa aquí la revisión de formato la reporta como
    # invasión en todas las presentaciones para público externo.
    pagina, con_cintilla = 1, set()
    titulos, con_nota, total = [], 0, 0
    firmas = {}
    display, seguidas, max_seguidas, por_seccion, actual = 0, 0, 0, {}, None

    for nivel, titulo, cont in laminas(cuerpo):
        pagina += 1
        if nivel == 1:
            actual = titulo
            por_seccion[actual] = 0
            continue
        if "\\cintilla" in cont:
            con_cintilla.add(pagina)
        total += 1
        if actual:
            por_seccion[actual] += 1
        if re.search(DISPLAY, cont):
            display += 1
            seguidas += 1
            max_seguidas = max(max_seguidas, seguidas)
        else:
            seguidas = 0
        vis, en_latex = texto_visible(cont)
        if "{.notes}" in cont or "\\note{" in cont:
            con_nota += 1
        n = palabras(vis, en_latex)
        if n > TOPE_PALABRAS and "tablaIrem" not in cont:
            fallos.append(f"«{titulo[:40]}»: {n} palabras (tope {TOPE_PALABRAS})")
        vin = [l for l in vis.split("\n") if re.match(r"^\s*[-*] ", l)]
        if len(vin) > TOPE_VINETAS:
            fallos.append(f"«{titulo[:40]}»: {len(vin)} viñetas (tope {TOPE_VINETAS})")
        for v in vin:
            if len(v.strip()) > TOPE_CARACTERES_VINETA:
                fallos.append(f"«{titulo[:40]}»: viñeta de más de dos renglones")
        if titulo and not titulo.startswith("{"):
            titulos.append(titulo)
            # La firma de una lámina es el título MÁS su primer renglón de
            # cuerpo. Repetir solo el título es la convención de la casa: el
            # título nombra la sección y se repite en todas las láminas de
            # esa sección; lo que distingue una de otra es el primer renglón.
            # Lo que no puede repetirse es el par completo.
            primeros = [l.strip() for l in vis.split("\n") if l.strip()]
            primero = re.sub(r"^[-*]\s+", "", primeros[0])[:60] if primeros else ""
            firmas.setdefault((titulo, primero), []).append(titulo)

    for (tit, primero), veces in firmas.items():
        if len(veces) > 1:
            detalle = f"primer renglón «{primero[:38]}»" if primero else "sin renglón de cuerpo que las distinga"
            fallos.append(
                f"{len(veces)} láminas indistinguibles con el título «{tit[:36]}» "
                f"({detalle}); califícalas con dos puntos")

    # Marcadores sin resolver. Se excluyen los enlaces de Markdown y los
    # argumentos opcionales de los comandos del formato, p. ej.
    # \laminaGracias[Preguntas], que no son marcadores.
    limpio = re.sub(r"\[[^\]]*\]\([^)]*\)", "", cuerpo)
    limpio = re.sub(r"\\[a-zA-Z]+\[[^\]]*\]", "", limpio)
    limpio = re.sub(r"```.*?```", lambda m: re.sub(r"\\[a-zA-Z]+\[[^\]]*\]", "", m.group(0)), limpio, flags=re.S)
    for m in re.findall(r"\[[^\]\n]{6,}\]", limpio):
        fallos.append(f"marcador sin resolver: {m}")

    if total and con_nota / total < 0.5:
        avisos.append(f"solo {con_nota} de {total} láminas tienen nota del presentador")

    # Portadillas que no se ganan su lugar.
    for nombre, n in por_seccion.items():
        if n < MIN_LAMINAS_POR_SECCION:
            fallos.append(
                f"«{nombre}» tiene {n} lámina de contenido y aun así lleva portadilla; "
                f"quítale el # o dale una segunda")

    # Láminas de una sola afirmación.
    if display and total // LAMINAS_POR_DISPLAY < display:
        fallos.append(
            f"{display} láminas de una sola frase para {total} de contenido "
            f"(máximo {max(1, total // LAMINAS_POR_DISPLAY)})")
    if max_seguidas > 1:
        fallos.append(f"{max_seguidas} láminas de una sola frase seguidas")

    # Cuánto del deck va casi en blanco: portada + portadillas + despliegue + cierre.
    paginas = 1 + len(secciones) + total
    vacias = len(secciones) + display + 1
    if paginas and vacias / paginas > TOPE_VACIAS:
        avisos.append(
            f"por estructura, {vacias} de {paginas} láminas llevarían una frase o menos; "
            f"la cuenta buena es la de tinta, más abajo")

    return secciones, total, con_nota, fallos, avisos, titulos, con_cintilla


def revisar_formato(pdf, con_cintilla=()):
    from PIL import Image

    # El fondo del formato no es blanco: lleva un degradado gris tenue que baja
    # hasta 227. Comparar contra blanco puro daría todas las láminas como
    # llenas de tinta. Así que se compara contra el fondo mismo: es tinta lo
    # que se aparta de él. Si no se encuentra el asset, se cae a un umbral
    # tolerante, que es peor pero no rompe.
    fondo = None
    for base in (pdf.parent, pdf.parent.parent):
        cand = base / "_extensions" / "irem" / "logos" / "fondo.png"
        if cand.exists():
            fondo = Image.open(cand).convert("RGB")
            break
    carpeta = pdf.parent / ".revision"
    carpeta.mkdir(exist_ok=True)
    for viejo in carpeta.glob("rev-*.png"):
        viejo.unlink()
    r = subprocess.run(
        ["Rscript", "-e",
         f'n <- pdftools::pdf_info("{pdf}")$pages;'
         f'pdftools::pdf_convert("{pdf}", pages=1:n, dpi=110, format="png",'
         f' filenames=file.path("{carpeta}", sprintf("rev-%02d.png", 1:n)))'],
        capture_output=True, text=True)
    pngs = sorted(carpeta.glob("rev-*.png"))
    if not pngs:
        return None, [f"no pude renderizar los PNG: {r.stderr.strip()[:200]}"]

    fallos = []
    flacas = []
    for i, f in enumerate(pngs, 1):
        im = Image.open(f).convert("RGB")
        W, H = im.size
        px = im.load()
        mmy, mmx = 90.0 / H, 160.0 / W
        if fondo is not None:
            fpx = fondo.resize((W, H), Image.BILINEAR).load()
            bl = lambda p, x=0, y=0: True  # se redefine abajo con coordenadas
            def bl(p, x=None, y=None, _f=fpx):
                if x is None:
                    return all(v > 225 for v in p)
                return max(abs(p[i] - _f[x, y][i]) for i in range(3)) < 14
        else:
            bl = lambda p, x=None, y=None: all(v > 225 for v in p)

        # Cuánta tinta lleva de verdad la lámina. Es la única manera honesta de
        # detectar la lámina de una sola frase: contar portadillas y comandos
        # de despliegue se queda corto, porque no ve la lámina de dos viñetas
        # sueltas que en pantalla también se lee como vacía. Se mide solo el
        # área de contenido, sin los logotipos ni la franja.
        y0, y1 = int(4 / mmy), int(74 / mmy)
        tinta = sum(1 for y in range(y0, y1) for x in range(W) if not bl(px[x, y], x, y))
        pct = 100.0 * tinta / ((y1 - y0) * W)
        if i > 1 and pct < TINTA_MINIMA:
            flacas.append((i, pct))
        if i == 1:
            arriba = sum(1 for c in range(int(0.60 * W), W) if not bl(px[c, 0], c, 0))
            der = sum(1 for y in range(0, int(0.93 * H)) if not bl(px[W - 1, y], W - 1, y))
            if arriba < 50:
                fallos.append(f"lámina 1: la foto no sangra por arriba")
            if der < 50:
                fallos.append(f"lámina 1: la foto no sangra por la derecha")
            continue
        z = sum(1 for y in range(int(74.0 / mmy), int(77.0 / mmy))
                for c in range(int(0.16 * W), int(0.88 * W)) if not bl(px[c, y], c, y))
        if z > 60 and i not in con_cintilla:
            fallos.append(f"lámina {i}: hay contenido en la banda de los logotipos")
        # Entre la base de los logotipos y la franja del pie no va nada. La
        # franja misma ocupa de 86 mm al borde, así que la ventana vigilada es
        # la de en medio; si se mirara hasta el borde, la propia franja daría
        # positivo en todas las láminas.
        b = sum(1 for y in range(int(83.5 / mmy), int(85.8 / mmy))
                for c in range(W) if not bl(px[c, y], c, y))
        if b > 60:
            fallos.append(f"lámina {i}: hay contenido pegado a la franja del pie")
        # Y la franja tiene que estar: azul sobre verde, de borde a borde.
        medio = W // 2
        azul = px[medio, int(87.0 / mmy)]
        verde = px[medio, int(89.0 / mmy)]
        cerca = lambda c, r, g, bb: abs(c[0]-r) < 26 and abs(c[1]-g) < 26 and abs(c[2]-bb) < 26
        if not cerca(azul, 0x36, 0x7C, 0xBC) or not cerca(verde, 0x98, 0xCE, 0x63):
            fallos.append(f"lámina {i}: falta la franja azul y verde del pie")
        # 150.5 y no 148.8: el recuadro \realce llega a 149.5 mm por diseño
        # del master, que le da margen derecho de 10.5 mm y no de 11.2.
        d = sum(1 for y in range(int(4.0 / mmy), int(74.0 / mmy))
                for c in range(int(150.5 / mmx), W) if not bl(px[c, y], c, y))
        if d > 40:
            fallos.append(f"lámina {i}: algo se sale por la derecha")
        iz = sum(1 for y in range(int(4.0 / mmy), int(74.0 / mmy))
                 for c in range(0, int(9.5 / mmx)) if not bl(px[c, y], c, y))
        if iz > 40:
            fallos.append(f"lámina {i}: algo se sale por la izquierda")
    return len(pngs), fallos


def revisar_pptx(pptx, qmd_laminas):
    """Revisa el .pptx sin renderizarlo.

    No sustituye mirar el archivo en PowerPoint, pero sí caza lo que se ve mal
    solo al proyectar: una forma metida en la banda de los logotipos, un texto
    que no cabe en su caja y, sobre todo, un renglón que se quedó con la
    tipografía del tema (Arial) o con el azul de Google Slides que arrastran los
    layouts del master. Eso último no se nota en la lámina hasta que alguien la
    abre al lado de otra."""
    from pptx import Presentation
    from pptx.util import Emu

    try:
        from PIL import ImageFont
        fuente = None
        for base in (Path.home() / "Library/Fonts", Path("/Library/Fonts")):
            if (base / "Montserrat-Regular.otf").exists():
                fuente = ImageFont.truetype(str(base / "Montserrat-Regular.otf"), 1000)
                break
        if fuente is None:
            r = subprocess.run(["kpsewhich", "Montserrat-Regular.otf"],
                               capture_output=True, text=True)
            if r.stdout.strip():
                fuente = ImageFont.truetype(r.stdout.strip(), 1000)
    except Exception:
        fuente = None

    def ancho_mm(t, pt):
        if fuente is None:
            return len(t) * 0.55 * pt * 25.4 / 72.0
        return fuente.getlength(t) / 1000.0 * pt * 25.4 / 72.0

    prs = Presentation(str(pptx))
    mm = lambda v: Emu(v).mm
    fallos, avisos = [], []
    con_nota = 0

    for i, lamina in enumerate(prs.slides, 1):
        if lamina.has_notes_slide and lamina.notes_slide.notes_text_frame.text.strip():
            con_nota += 1
        for sh in lamina.shapes:
            nombre = getattr(sh, "name", "forma")
            # Las piezas del pie y del cierre viven ahí a propósito: el
            # generador las nombra para poder saltárselas aquí.
            if nombre.startswith(("pie-", "cierre-")):
                continue
            # Las imágenes del pie (logotipos, cintillas, banda del cierre) van
            # ahí a propósito; lo que no puede bajar es el texto y las tablas.
            if sh.shape_type == 13:
                continue
            try:
                x, y = mm(sh.left), mm(sh.top)
                w, h = mm(sh.width), mm(sh.height)
            except (TypeError, ValueError):
                continue
            if sh.has_text_frame and not sh.text_frame.text.strip():
                continue
            if y + h > AREA_FIN_PPT + 1 and sh.has_text_frame:
                # Una caja de texto alta pero medio vacía no es un problema: lo
                # que importa es dónde termina el TEXTO, no la caja.
                pt = 21
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            pt = r.font.size.pt
                        break
                    break
                lineas = max(1, len(sh.text_frame.text.split("\n")))
                alto_texto = lineas * pt * ALTURA_LINEA * 1.2 * 25.4 / 72.0
                if y + alto_texto > AREA_FIN_PPT + 1:
                    fallos.append(f"lámina {i}: «{nombre}» baja hasta "
                                  f"{y + alto_texto:.0f} mm, dentro de la banda de "
                                  f"los logotipos (el área útil termina en "
                                  f"{AREA_FIN_PPT:.0f})")
            # La portada tiene su propia geometría (el título arranca en
            # 13.9 mm, no en el margen del cuerpo), así que los márgenes del
            # área de contenido no le aplican. Es la misma excepción que hace
            # la revisión del PDF.
            if i == 1:
                pass
            elif x < MARGEN_IZQ_PPT - 1:
                fallos.append(f"lámina {i}: «{nombre}» empieza en {x:.0f} mm, "
                              f"fuera del margen izquierdo")
            elif x + w > MARGEN_DER_PPT + 1:
                fallos.append(f"lámina {i}: «{nombre}» llega a {x + w:.0f} mm, "
                              f"fuera del margen derecho")
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if not r.text.strip():
                        continue
                    if r.font.name not in ("Montserrat", "Menlo"):
                        fallos.append(f"lámina {i}: «{r.text[:28]}» quedó en "
                                      f"{r.font.name or 'la fuente del tema'}, "
                                      f"no en Montserrat")
                    if r.font.size is None:
                        fallos.append(f"lámina {i}: «{r.text[:28]}» sin tamaño "
                                      f"propio; hereda el del layout")
            # Las cajas de tamaño fijo del formato (numerales, conceptos,
            # realce) no crecen: si el texto no cabe, se sale por abajo.
            if sh.shape_type == 1 and sh.has_text_frame:
                tf = sh.text_frame
                pt = 19
                for p in tf.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            pt = max(pt, r.font.size.pt)
                util = w - mm(tf.margin_left or 0) - mm(tf.margin_right or 0)
                alto = 0.0
                for p in tf.paragraphs:
                    texto = "".join(r.text for r in p.runs)
                    if not texto:
                        continue
                    tam = next((r.font.size.pt for r in p.runs if r.font.size), pt)
                    n, linea = 1, ""
                    for palabra in texto.split():
                        prueba = f"{linea} {palabra}".strip()
                        if ancho_mm(prueba, tam) > util and linea:
                            n += 1
                            linea = palabra
                        else:
                            linea = prueba
                    alto += n * tam * ALTURA_LINEA * 25.4 / 72.0
                if alto > h - mm(tf.margin_top or 0) - mm(tf.margin_bottom or 0):
                    fallos.append(f"lámina {i}: el texto «{tf.text[:28]}» no cabe "
                                  f"en su caja; acórtalo")

    n = len(prs.slides._sldIdLst)
    if qmd_laminas and n != qmd_laminas:
        fallos.append(f"el .pptx tiene {n} láminas y el .qmd describe {qmd_laminas}")
    if n and con_nota / n < 0.5:
        avisos.append(f"solo {con_nota} de {n} láminas llevan nota en el panel "
                      f"de notas de PowerPoint")
    return n, fallos, avisos


def main():
    if len(sys.argv) != 2:
        sys.exit("Uso: ./revisar.py <archivo>.qmd")
    qmd = Path(sys.argv[1]).resolve()
    if not qmd.exists():
        sys.exit(f"No existe: {qmd}")

    secciones, total, con_nota, fallos_c, avisos, titulos, con_cintilla = revisar_contenido(qmd)

    # La espina puede ir en portadillas (#) o en los títulos de lámina (##),
    # que es lo habitual. En el segundo caso el título es el nombre de la
    # sección, quizá calificado: «Lo que vemos: flujo», «Lo que proponemos
    # para puntos urbanos».
    def raiz(t):
        for e in ESPINA:
            if t == e or t.startswith(e + ":") or t.startswith(e + " "):
                return e
        return None
    en_titulos, visto = [], None
    for t in titulos:
        r = raiz(t)
        if r and r != visto:
            en_titulos.append(r); visto = r
    # Si las portadillas llevan los nombres de la espina, mandan ellas; si no,
    # se buscan en los títulos, que es la convención por omisión.
    if set(ESPINA) & set(secciones):
        espina, donde = secciones, "portadillas"
    elif en_titulos:
        espina, donde = en_titulos, "títulos"
    else:
        espina, donde = secciones, "portadillas"
    externa = bool(set(ESPINA) & set(espina))

    print(f"\n{qmd.name}   {total} láminas, {con_nota} con nota")
    print(f"espina (en {donde}): {' → '.join(espina) if espina else '(ninguna)'}")

    if externa and espina != ESPINA:
        fallos_c.append(
            "parece una propuesta a un externo pero la espina no coincide.\n"
            f"    esperada: {' → '.join(ESPINA)}")

    print("\nCONTENIDO")
    for f in fallos_c:
        print(f"  !! {f}")
    for a in avisos:
        print(f"  ·  {a}")
    if not fallos_c and not avisos:
        print("  todo en regla")

    # Cuántas láminas TIENE que haber: portada, una por sección con portadilla
    # y una por lámina de contenido. Se compara contra lo que salió compilado,
    # porque las dos formas de que sobre una lámina (un comentario de HTML
    # antes de la primera, un bloque suelto) no dan ningún error y no se ven
    # hasta proyectar.
    esperadas = 1 + len(secciones) + total

    print("\nFORMATO")
    pdf = qmd.with_suffix(".pdf")
    fallos_f = []
    if not pdf.exists():
        print(f"  no hay {pdf.name}; compila primero con ./renderizar.sh")
        n_fmt = None
    else:
        n_fmt, fallos_f = revisar_formato(pdf, con_cintilla)
        if n_fmt and n_fmt != esperadas:
            fallos_f.insert(0, f"el PDF tiene {n_fmt} láminas y el .qmd describe "
                               f"{esperadas}. Si sobra una y sale en blanco, casi "
                               f"siempre es un comentario de HTML entre el "
                               f"encabezado y la primera lámina: pásalo al "
                               f"encabezado, con # delante")
        for f in fallos_f:
            print(f"  !! {f}")
        if not fallos_f:
            print(f"  {n_fmt} láminas, nada fuera de caja")

    fallos_p = []
    pptx = qmd.with_suffix(".pptx")
    if pptx.exists():
        print("\nPOWERPOINT")
        n_ppt, fallos_p, avisos_p = revisar_pptx(pptx, esperadas)
        for f in fallos_p:
            print(f"  !! {f}")
        for a in avisos_p:
            print(f"  ·  {a}")
        if not fallos_p and not avisos_p:
            print(f"  {n_ppt} láminas, nada fuera de caja")
        print("  (esto no lo renderiza: ábrelo en PowerPoint y míralo, o pasa")
        print("   ./pptx-a-pdf.sh para convertirlo y revisar los PNG)")

    print("\nEsto no reemplaza mirar las láminas: revisa que el argumento se")
    print("entienda y que ninguna lámina se lea en voz alta.")
    sys.exit(1 if (fallos_c or fallos_f or fallos_p) else 0)


if __name__ == "__main__":
    main()
